#!/usr/bin/env python3
"""生成 DLQI 中期汇报 PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色定义
BG = RGBColor(0x0a, 0x0a, 0x0a)
CARD_BG = RGBColor(0x14, 0x14, 0x14)
WHITE = RGBColor(0xe5, 0xe7, 0xeb)
DIM = RGBColor(0x9c, 0xa3, 0xaf)
BLUE = RGBColor(0x3b, 0x82, 0xf6)
GREEN = RGBColor(0x10, 0xb9, 0x81)
RED = RGBColor(0xef, 0x44, 0x44)
YELLOW = RGBColor(0xf5, 0x9e, 0x0b)


def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=CARD_BG, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    p.font.name = "Microsoft YaHei"
    return txBox


def add_bullet_text(slide, left, top, width, height, items, size=13, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(6)
    return txBox


# ==================== Slide 1: 封面 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1),
         "DLQI — 深度学习量化交易研究平台", size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.6),
         "Deep Learning Quantitative Intelligence", size=20, color=BLUE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.5),
         "中期汇报", size=24, color=DIM, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
         "2026年4月", size=16, color=DIM, align=PP_ALIGN.CENTER)

# ==================== Slide 2: 项目概述 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
         "项目概述", size=28, color=WHITE, bold=True)

# 研究目标
add_shape(slide, Inches(0.8), Inches(1.2), Inches(5.8), Inches(2.5))
add_text(slide, Inches(1.1), Inches(1.3), Inches(5.3), Inches(0.4),
         "研究目标", size=16, color=BLUE, bold=True)
add_bullet_text(slide, Inches(1.1), Inches(1.8), Inches(5.3), Inches(1.8), [
    "使用深度学习模型预测美股价格走势",
    "对比 Transformer、LSTM、LightGBM、XGBoost 四种模型",
    "研究对象：AAPL、AMZN、GOOGL、MSFT、NVDA 五只美股",
    "构建完整的量化交易回测与风险评估体系",
], size=13, color=WHITE)

# 技术栈
add_shape(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(2.5))
add_text(slide, Inches(7.2), Inches(1.3), Inches(5.3), Inches(0.4),
         "技术栈", size=16, color=BLUE, bold=True)
add_bullet_text(slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(1.8), [
    "后端：FastAPI + SQLAlchemy + Supabase PostgreSQL",
    "前端：Vue 3 + TypeScript + ECharts + Tailwind CSS",
    "模型：PyTorch (LSTM/Transformer) + LightGBM + XGBoost",
    "训练：Google Colab T4 GPU 远程训练",
], size=13, color=WHITE)

# 系统架构
add_shape(slide, Inches(0.8), Inches(4.0), Inches(11.9), Inches(3.0))
add_text(slide, Inches(1.1), Inches(4.1), Inches(5), Inches(0.4),
         "系统架构", size=16, color=BLUE, bold=True)
add_bullet_text(slide, Inches(1.1), Inches(4.6), Inches(5.5), Inches(2.2), [
    "8 个功能页面：仪表盘、策略、模型、回测、交易、风险、数据、调参",
    "6 组 RESTful API：数据、模型、回测、风险、交易、任务",
    "Google Drive 任务队列实现本地与 Colab 协同训练",
], size=13, color=WHITE)
add_bullet_text(slide, Inches(7.0), Inches(4.6), Inches(5.5), Inches(2.2), [
    "S&P 500 全部 503 只成分股 10 年日线数据",
    "12 维特征：OHLCV + 收益率 + MA + 波动率 + RSI + MACD",
    "回测引擎：100 万初始资金，含手续费和滑点",
], size=13, color=WHITE)


# ==================== Slide 3: 模型对比 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "模型性能对比", size=28, color=WHITE, bold=True)

# 表格标题行
headers = ["模型类型", "平均 Sharpe", "平均收益率", "方向准确率", "平均最大回撤", "模型数"]
col_widths = [2.2, 1.8, 1.8, 1.8, 1.8, 1.2]
x_start = 0.8

# Header
y = 1.3
for j, (h, w) in enumerate(zip(headers, col_widths)):
    add_shape(slide, Inches(x_start + sum(col_widths[:j])), Inches(y), Inches(w), Inches(0.45), fill_color=RGBColor(0x1e, 0x1e, 0x1e))
    add_text(slide, Inches(x_start + sum(col_widths[:j]) + 0.1), Inches(y + 0.05), Inches(w - 0.2), Inches(0.35),
             h, size=12, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# Data rows
data = [
    ["LightGBM", "0.89", "23.8%", "53.6%", "-21.8%", "5"],
    ["Transformer (分类)", "0.52", "5.5%", "49.7%", "-14.2%", "5"],
    ["LSTM", "0.22", "7.2%", "49.9%", "-13.8%", "5"],
    ["XGBoost", "-0.20", "-0.4%", "50.5%", "-21.5%", "5"],
]
row_colors = [GREEN, BLUE, WHITE, RED]

for i, (row, rc) in enumerate(zip(data, row_colors)):
    y = 1.8 + i * 0.5
    for j, (val, w) in enumerate(zip(row, col_widths)):
        bg = RGBColor(0x12, 0x12, 0x12) if i % 2 == 0 else CARD_BG
        add_shape(slide, Inches(x_start + sum(col_widths[:j])), Inches(y), Inches(w), Inches(0.45), fill_color=bg)
        c = rc if j == 0 or j == 1 else WHITE
        add_text(slide, Inches(x_start + sum(col_widths[:j]) + 0.1), Inches(y + 0.05), Inches(w - 0.2), Inches(0.35),
                 val, size=12, color=c, bold=(i == 0), align=PP_ALIGN.CENTER)

# 关键发现
add_shape(slide, Inches(0.8), Inches(4.2), Inches(5.8), Inches(2.8))
add_text(slide, Inches(1.1), Inches(4.3), Inches(5.3), Inches(0.4),
         "关键发现", size=16, color=GREEN, bold=True)
add_bullet_text(slide, Inches(1.1), Inches(4.8), Inches(5.3), Inches(2.0), [
    "LightGBM 回测表现最优 (平均 Sharpe 0.89)",
    "Transformer 分类版排第二 (Sharpe 0.52)，回撤最小 (-14.2%)",
    "回归版 Transformer 存在模型坍缩问题，改为分类任务后解决",
    "所有模型超额收益为负，市场 beta 贡献了大部分收益",
], size=13, color=WHITE)

add_shape(slide, Inches(6.9), Inches(4.2), Inches(5.8), Inches(2.8))
add_text(slide, Inches(7.2), Inches(4.3), Inches(5.3), Inches(0.4),
         "Top 5 模型（含超额收益）", size=16, color=GREEN, bold=True)
add_bullet_text(slide, Inches(7.2), Inches(4.8), Inches(5.3), Inches(2.0), [
    "#1  Transformer GOOGL — Sharpe 3.38, 超额 -46.3%",
    "#2  LightGBM GOOGL — Sharpe 3.29, 超额 -4.8%",
    "#3  XGBoost GOOGL — Sharpe 1.62, 超额 -48.8%",
    "#4  LightGBM AAPL — Sharpe 1.36, 超额 -1.4%",
    "#5  LightGBM NVDA — Sharpe 1.21, 超额 -37.0%",
], size=13, color=WHITE)


# ==================== Slide 4: Transformer 改进 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "核心工作：Transformer 架构改进", size=28, color=WHITE, bold=True)

# 改进前
add_shape(slide, Inches(0.8), Inches(1.2), Inches(5.8), Inches(2.5), border_color=RED)
add_text(slide, Inches(1.1), Inches(1.3), Inches(5.3), Inches(0.4),
         "改进前", size=16, color=RED, bold=True)
add_bullet_text(slide, Inches(1.1), Inches(1.8), Inches(5.3), Inches(1.8), [
    "随机可学习位置编码 — 小数据集难以收敛",
    "无因果掩码 — 训练时泄露未来信息",
    "模型过大：d_model=128, 4层, 8头 (~50万参数)",
    "单股票训练：仅 ~1200 条数据，严重过拟合",
    "结果：平均 Sharpe 0.83，方向准确率 48.2%",
], size=13, color=WHITE)

# 改进后
add_shape(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(2.5), border_color=GREEN)
add_text(slide, Inches(7.2), Inches(1.3), Inches(5.3), Inches(0.4),
         "改进后", size=16, color=GREEN, bold=True)
add_bullet_text(slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(1.8), [
    "正弦位置编码 — 无需学习，小数据更稳定",
    "因果掩码 — 训练/推理分布一致",
    "模型缩小：d_model=64, 2层, 4头 (~10万参数)",
    "改为分类任务（涨/跌），解决回归坍缩问题",
    "结果：方向准确率 52.0%（超过随机 50%）",
], size=13, color=WHITE)

# 改进效果数字
add_shape(slide, Inches(0.8), Inches(4.0), Inches(3.8), Inches(1.5))
add_text(slide, Inches(0.8), Inches(4.1), Inches(3.8), Inches(0.4),
         "任务类型", size=14, color=DIM, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.8), Inches(4.5), Inches(3.8), Inches(0.6),
         "回归 → 分类", size=24, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(4.8), Inches(4.0), Inches(3.8), Inches(1.5))
add_text(slide, Inches(4.8), Inches(4.1), Inches(3.8), Inches(0.4),
         "方向准确率", size=14, color=DIM, align=PP_ALIGN.CENTER)
add_text(slide, Inches(4.8), Inches(4.5), Inches(3.8), Inches(0.6),
         "48.2% → 52.0%  (+3.8pp)", size=24, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(8.8), Inches(4.0), Inches(3.8), Inches(1.5))
add_text(slide, Inches(8.8), Inches(4.1), Inches(3.8), Inches(0.4),
         "模型坍缩", size=14, color=DIM, align=PP_ALIGN.CENTER)
add_text(slide, Inches(8.8), Inches(4.5), Inches(3.8), Inches(0.6),
         "已解决", size=24, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# 训练优化
add_shape(slide, Inches(0.8), Inches(5.8), Inches(11.9), Inches(1.3))
add_text(slide, Inches(1.1), Inches(5.9), Inches(5), Inches(0.4),
         "训练优化", size=16, color=BLUE, bold=True)
add_bullet_text(slide, Inches(1.1), Inches(6.3), Inches(5.3), Inches(0.7), [
    "AdamW + Cosine Annealing LR + 梯度裁剪",
    "Early Stopping (patience=10)",
], size=12, color=WHITE)
add_bullet_text(slide, Inches(7.0), Inches(6.3), Inches(5.3), Inches(0.7), [
    "LazySeqDataset 按需生成序列，避免 OOM",
    "逐股票构建序列，避免跨股票边界污染",
], size=12, color=WHITE)


# ==================== Slide 5: 数据与特征 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "数据与特征工程", size=28, color=WHITE, bold=True)

# 数据规模
add_shape(slide, Inches(0.8), Inches(1.2), Inches(3.8), Inches(2.0))
add_text(slide, Inches(0.8), Inches(1.3), Inches(3.8), Inches(0.4),
         "数据规模", size=14, color=DIM, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.8), Inches(1.7), Inches(3.8), Inches(0.5),
         "503 只股票", size=28, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.8), Inches(2.3), Inches(3.8), Inches(0.4),
         "S&P 500 全部成分股 × 10 年日线", size=12, color=DIM, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(4.8), Inches(1.2), Inches(3.8), Inches(2.0))
add_text(slide, Inches(4.8), Inches(1.3), Inches(3.8), Inches(0.4),
         "总数据量", size=14, color=DIM, align=PP_ALIGN.CENTER)
add_text(slide, Inches(4.8), Inches(1.7), Inches(3.8), Inches(0.5),
         "1,245,066 条", size=28, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(4.8), Inches(2.3), Inches(3.8), Inches(0.4),
         "训练集 841,217 / 验证集 156,546", size=12, color=DIM, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(8.8), Inches(1.2), Inches(3.8), Inches(2.0))
add_text(slide, Inches(8.8), Inches(1.3), Inches(3.8), Inches(0.4),
         "数据源", size=14, color=DIM, align=PP_ALIGN.CENTER)
add_text(slide, Inches(8.8), Inches(1.7), Inches(3.8), Inches(0.5),
         "yfinance", size=28, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(8.8), Inches(2.3), Inches(3.8), Inches(0.4),
         "Yahoo Finance API，免费开源", size=12, color=DIM, align=PP_ALIGN.CENTER)

# 特征列表
add_shape(slide, Inches(0.8), Inches(3.5), Inches(11.9), Inches(3.5))
add_text(slide, Inches(1.1), Inches(3.6), Inches(5), Inches(0.4),
         "12 维特征", size=16, color=BLUE, bold=True)

features = [
    ("基础价格", "Open, High, Low, Close, Volume"),
    ("收益率", "日收益率 (pct_change)"),
    ("均线", "MA5 (5日均线), MA20 (20日均线)"),
    ("波动率", "20日滚动标准差"),
    ("动量", "RSI (14日相对强弱指标)"),
    ("趋势", "MACD + MACD Signal (12/26/9)"),
]
for i, (name, desc) in enumerate(features):
    y_pos = 4.1 + i * 0.42
    x_pos = 1.1 if i < 3 else 7.0
    y_pos = 4.1 + (i % 3) * 0.42
    add_text(slide, Inches(x_pos), Inches(y_pos), Inches(1.5), Inches(0.35),
             name, size=12, color=GREEN, bold=True)
    add_text(slide, Inches(x_pos + 1.5), Inches(y_pos), Inches(4), Inches(0.35),
             desc, size=12, color=WHITE)

add_text(slide, Inches(1.1), Inches(5.5), Inches(10), Inches(0.8),
         "序列长度：60 个交易日（约 3 个月）  |  数据分割：70% 训练 / 15% 验证 / 15% 测试  |  标准化：StandardScaler",
         size=12, color=DIM)


# ==================== Slide 6: 系统功能 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "系统功能展示", size=28, color=WHITE, bold=True)

pages = [
    ("仪表盘", "模型排名、KPI 概览、策略推荐"),
    ("模型管理", "模型列表、特征重要性、性能对比"),
    ("回测分析", "回测结果表、资金曲线、相关性矩阵"),
    ("模拟交易", "虚拟投资组合、历史模拟、交易记录"),
    ("风险监控", "VaR 计算、压力测试、风险预警"),
    ("数据管理", "数据同步、质量统计、批量下载"),
    ("参数调优", "训练配置、进度监控、结果刷新"),
    ("策略推荐", "复合评分、一键部署、权益曲线"),
]

for i, (name, desc) in enumerate(pages):
    col = i % 4
    row = i // 4
    x = 0.8 + col * 3.1
    y = 1.2 + row * 2.8
    add_shape(slide, Inches(x), Inches(y), Inches(2.9), Inches(2.4))
    add_text(slide, Inches(x + 0.2), Inches(y + 0.2), Inches(2.5), Inches(0.4),
             name, size=16, color=BLUE, bold=True)
    add_text(slide, Inches(x + 0.2), Inches(y + 0.7), Inches(2.5), Inches(1.5),
             desc, size=12, color=DIM)


# ==================== Slide 7: 后期计划 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "后期工作计划", size=28, color=WHITE, bold=True)

# 模型层面
add_shape(slide, Inches(0.8), Inches(1.2), Inches(3.7), Inches(5.5))
add_text(slide, Inches(1.1), Inches(1.3), Inches(3.3), Inches(0.4),
         "模型优化", size=16, color=BLUE, bold=True)
add_bullet_text(slide, Inches(1.1), Inches(1.8), Inches(3.3), Inches(4.5), [
    "特征增强：引入布林带、ATR、OBV 等更多技术指标",
    "集成学习：Transformer + LightGBM 信号融合",
    "超参数优化：Optuna 系统搜索最优配置",
    "注意力可视化：分析模型关注的时间步",
], size=12, color=WHITE)

# 策略层面
add_shape(slide, Inches(4.7), Inches(1.2), Inches(3.7), Inches(5.5))
add_text(slide, Inches(5.0), Inches(1.3), Inches(3.3), Inches(0.4),
         "策略改进", size=16, color=GREEN, bold=True)
add_bullet_text(slide, Inches(5.0), Inches(1.8), Inches(3.3), Inches(4.5), [
    "动态仓位：根据模型置信度调整仓位比例",
    "止损止盈：引入风险控制规则",
    "组合优化：多股票权重优化 (Markowitz)",
    "交易成本建模：更精确的滑点和冲击成本",
], size=12, color=WHITE)

# 论文层面
add_shape(slide, Inches(8.6), Inches(1.2), Inches(3.7), Inches(5.5))
add_text(slide, Inches(8.9), Inches(1.3), Inches(3.3), Inches(0.4),
         "论文撰写", size=16, color=YELLOW, bold=True)
add_bullet_text(slide, Inches(8.9), Inches(1.8), Inches(3.3), Inches(4.5), [
    "整理实验数据和对比分析",
    "撰写模型实验章节",
    "绘制架构图和流程图",
    "完成毕业论文初稿",
], size=12, color=WHITE)


# ==================== Slide 8: 总结 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.8),
         "总结", size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(2), Inches(2.5), Inches(9), Inches(3.5))
add_bullet_text(slide, Inches(2.5), Inches(2.8), Inches(8), Inches(3.0), [
    "完成了完整的量化交易研究平台（前端 + 后端 + 训练 + 回测）",
    "实现 4 种模型对比，LightGBM 回测表现最优 (Sharpe 0.89)",
    "Transformer 改为分类任务后方向准确率 52%，解决了回归坍缩问题",
    "发现模型超额收益有限，市场方向预测本质上是极低信噪比问题",
    "后续重点：特征增强、集成学习、更长训练、论文撰写",
], size=16, color=WHITE)

add_text(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5),
         "谢谢！", size=24, color=BLUE, bold=True, align=PP_ALIGN.CENTER)


# ==================== 保存 ====================
output_path = "/home/awan/DLQI/docs/DLQI_中期汇报.pptx"
prs.save(output_path)
print(f"PPT 已生成: {output_path}")
